"""
Joining Hands - Investor Pitch Deck
Design inspired by the Engineering Review reference deck.
White content slides, deep navy cover, clean corporate layout.
No GramConnect. Compact. Sophisticated.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# --- DESIGN TOKENS (Matched to reference deck) ---
class C:
    NAVY       = RGBColor(30, 39, 97)      # #1E2761 — cover/accent bg
    NAVY_MID   = RGBColor(46, 59, 130)     # #2E3B82 — icon circles
    LAVENDER   = RGBColor(202, 220, 252)   # #CADCFC — accent bars
    PANEL      = RGBColor(245, 247, 252)   # #F5F7FC — light panels
    WHITE      = RGBColor(255, 255, 255)
    BLACK      = RGBColor(30, 30, 40)      # Near-black text
    
    T1 = RGBColor(30, 30, 50)             # Headings
    T2 = RGBColor(80, 85, 105)            # Body text
    T3 = RGBColor(130, 135, 155)          # Muted / footnotes
    
    # Accent colors for cards
    SAFFRON    = RGBColor(255, 119, 0)     # #FF7700
    TEAL       = RGBColor(16, 185, 129)    # Success green
    BLUE       = RGBColor(59, 130, 246)    # Bright blue
    RED        = RGBColor(220, 53, 69)     # Alert/critical
    
    # White-slide text for dark panels
    W1 = RGBColor(255, 255, 255)
    W2 = RGBColor(200, 210, 235)
    W3 = RGBColor(150, 160, 190)
    
    FONT = 'Calibri'


def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- Helpers ---
    def set_bg(s, color):
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = color

    def rect(s, l, t, w, h, color, border=None, bw=Pt(1)):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        sh.fill.solid(); sh.fill.fore_color.rgb = color
        if border: sh.line.color.rgb = border; sh.line.width = bw
        else: sh.line.fill.background()
        return sh

    def rnd(s, l, t, w, h, color, border=None, bw=Pt(1)):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        sh.fill.solid(); sh.fill.fore_color.rgb = color
        if border: sh.line.color.rgb = border; sh.line.width = bw
        else: sh.line.fill.background()
        return sh

    def circ(s, l, t, sz, color):
        sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
        sh.fill.solid(); sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        return sh

    def tx(s, l, t, w, h, text, sz=14, color=C.T1, bold=False, italic=False, align=PP_ALIGN.LEFT):
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text; p.font.name = C.FONT; p.font.size = Pt(sz)
        p.font.color.rgb = color; p.font.bold = bold; p.font.italic = italic
        p.alignment = align
        return tf

    def mtx(s, l, t, w, h, lines):
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ln.get('text', '')
            p.font.name = C.FONT
            p.font.size = Pt(ln.get('sz', 13))
            p.font.color.rgb = ln.get('color', C.T2)
            p.font.bold = ln.get('bold', False)
            p.font.italic = ln.get('italic', False)
            p.alignment = ln.get('align', PP_ALIGN.LEFT)
            p.space_after = Pt(ln.get('sp', 4))

    def white_footer(s, num):
        tx(s, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
           "The Group of Joining Hands  |  Investor Pitch Deck", sz=9, color=C.T3)
        tx(s, Inches(11.8), Inches(7.05), Inches(1.0), Inches(0.3),
           str(num), sz=9, color=C.T3, bold=True, align=PP_ALIGN.RIGHT)

    def navy_footer(s, num):
        tx(s, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
           "The Group of Joining Hands  |  Investor Pitch Deck", sz=9, color=C.W3)
        tx(s, Inches(11.8), Inches(7.05), Inches(1.0), Inches(0.3),
           str(num), sz=9, color=C.W3, bold=True, align=PP_ALIGN.RIGHT)

    def section_head(s, tag, title):
        """White slide section header: small tag + large heading."""
        tx(s, Inches(0.6), Inches(0.45), Inches(8), Inches(0.3),
           tag, sz=11, color=C.NAVY, bold=True)
        tx(s, Inches(0.6), Inches(0.75), Inches(11), Inches(0.65),
           title, sz=30, color=C.T1, bold=True)

    def accent_dot(s, l, t, h=Inches(0.6)):
        """Small lavender accent bar (vertical) like in reference."""
        rect(s, l, t, Inches(0.07), h, C.LAVENDER)

    def feature_row(s, l, t, w, title, desc):
        """Accent-dot + title + description (reference style)."""
        accent_dot(s, l, t, Inches(0.55))
        tx(s, l + Inches(0.2), t, w - Inches(0.2), Inches(0.28),
           title, sz=14, color=C.T1, bold=True)
        tx(s, l + Inches(0.2), t + Inches(0.3), w - Inches(0.2), Inches(0.3),
           desc, sz=11, color=C.T2)

    def navy_card(s, l, t, w, h, title, lines_below):
        """Dark navy card with icon circle + title + sub-lines."""
        rect(s, l, t, w, h, C.NAVY)
        circ(s, l + (w - Inches(0.5)) / 2, t + Inches(0.25), Inches(0.5), C.NAVY_MID)
        tx(s, l, t + Inches(0.9), w, Inches(0.3),
           title, sz=13, color=C.W1, bold=True, align=PP_ALIGN.CENTER)
        mtx(s, l + Inches(0.15), t + Inches(1.25), w - Inches(0.3), h - Inches(1.4),
            [{'text': ln, 'sz': 10, 'color': C.W2, 'sp': 3} for ln in lines_below])

    # ===================================================================
    #  SLIDE 1 — COVER (Navy background, matching reference slide 1)
    # ===================================================================
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, C.NAVY)

    # Horizontal divider line
    rect(s1, Inches(0.6), Inches(3.05), Inches(12.1), Inches(0.02), C.LAVENDER)

    # Top tag
    tx(s1, Inches(0.6), Inches(1.8), Inches(8), Inches(0.35),
       "INVESTOR PITCH DECK", sz=13, color=C.LAVENDER, bold=True)

    # Main title
    tx(s1, Inches(0.6), Inches(3.4), Inches(11), Inches(1.0),
       "The Group of Joining Hands", sz=46, color=C.W1, bold=True)

    # Subtitle
    tx(s1, Inches(0.6), Inches(4.5), Inches(10), Inches(0.5),
       "Community  |  Career Network  |  Commute Platform", sz=16, color=C.W2)

    # Divider
    rect(s1, Inches(0.6), Inches(5.2), Inches(2.5), Inches(0.015), C.LAVENDER)

    # Bottom bar with portal names
    rect(s1, Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.4), C.NAVY_MID)
    tx(s1, Inches(0.6), Inches(5.65), Inches(12.1), Inches(0.3),
       "ProConnect  |  H Rides  |  Real-Time Messaging  |  Admin Control Center",
       sz=11, color=C.LAVENDER, bold=True, align=PP_ALIGN.CENTER)

    navy_footer(s1, 1)

    # ===================================================================
    #  SLIDE 2 — PROJECT OVERVIEW (White bg, reference slide 2 style)
    # ===================================================================
    s2 = prs.slides.add_slide(blank)
    set_bg(s2, C.WHITE)
    section_head(s2, "WHAT WE'VE BUILT", "Project Overview")

    # Left column — description + feature list
    mtx(s2, Inches(0.6), Inches(1.7), Inches(5.5), Inches(1.8), [
        {'text': 'Unified Community & Professional Networking Ecosystem', 'sz': 16, 'color': C.T1, 'bold': True, 'sp': 10},
        {'text': 'A full-stack platform that merges career networking and ride-hailing '
                 'under a single verified identity. Zero ads, zero trackers, zero framework '
                 'dependencies. Everything hand-coded for maximum performance.',
         'sz': 12, 'color': C.T2, 'sp': 6},
    ])

    # Feature rows (with lavender accent dots like reference)
    feature_row(s2, Inches(0.6), Inches(3.7), Inches(5.3),
                "ProConnect", "Professional network: feed, messaging, analytics, certifications")
    feature_row(s2, Inches(0.6), Inches(4.5), Inches(5.3),
                "H Rides", "Transit platform: GPS tracking, rider-captain chat, OTP verification")
    feature_row(s2, Inches(0.6), Inches(5.3), Inches(5.3),
                "Messaging & Admin", "Real-time chat, notifications, and full admin control panel")

    # Right panel — Tech Stack (light gray box like reference)
    rect(s2, Inches(6.6), Inches(1.7), Inches(6.1), Inches(4.85), C.PANEL)
    tx(s2, Inches(7.0), Inches(1.95), Inches(5.3), Inches(0.3),
       "TECH STACK", sz=13, color=C.NAVY, bold=True)

    tech_items = [
        ("Python http.server", "Lightweight REST API backend, no framework dependency"),
        ("SQLite", "Single-file relational database, zero configuration"),
        ("Vanilla JS / HTML / CSS", "Framework-free frontend, fully client-rendered SPA"),
    ]

    for i, (title, desc) in enumerate(tech_items):
        y = Inches(2.6) + i * Inches(1.35)
        circ(s2, Inches(7.0), y, Inches(0.5), C.NAVY)
        tx(s2, Inches(7.7), y - Inches(0.02), Inches(4.7), Inches(0.28),
           title, sz=14, color=C.T1, bold=True)
        tx(s2, Inches(7.7), y + Inches(0.3), Inches(4.7), Inches(0.45),
           desc, sz=11, color=C.T2)

    white_footer(s2, 2)

    # ===================================================================
    #  SLIDE 3 — SYSTEM ARCHITECTURE (White bg, reference slide 3 style)
    # ===================================================================
    s3 = prs.slides.add_slide(blank)
    set_bg(s3, C.WHITE)
    section_head(s3, "HOW IT FITS TOGETHER", "System Architecture")

    # Three navy architecture boxes (Client → Server → Data)
    arch_boxes = [
        ("Client", ["Browser (SPA)", "HTML / CSS / JavaScript"], Inches(0.8)),
        ("Application Server", ["Python http.server", "REST API + JWT Auth"], Inches(4.8)),
        ("Data Layer", ["SQLite database.db", "users, posts, rides, messages"], Inches(8.8)),
    ]

    for (title, sub, x) in arch_boxes:
        rect(s3, x, Inches(2.2), Inches(3.2), Inches(2.0), C.NAVY)
        circ(s3, x + Inches(1.2), Inches(2.4), Inches(0.6), C.NAVY_MID)
        tx(s3, x, Inches(3.15), Inches(3.2), Inches(0.3),
           title, sz=13, color=C.W1, bold=True, align=PP_ALIGN.CENTER)
        mtx(s3, x + Inches(0.15), Inches(3.5), Inches(2.9), Inches(0.7),
            [{'text': s, 'sz': 10, 'color': C.W2, 'sp': 2, 'align': PP_ALIGN.CENTER} for s in sub])

    # Connector arrows (lavender rectangles like reference)
    rect(s3, Inches(4.0), Inches(3.05), Inches(0.8), Inches(0.2), C.LAVENDER)
    rect(s3, Inches(8.0), Inches(3.05), Inches(0.8), Inches(0.2), C.LAVENDER)

    # Bottom info panel
    rect(s3, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.1), C.PANEL)
    circ(s3, Inches(1.1), Inches(4.95), Inches(0.45), C.NAVY)
    mtx(s3, Inches(1.8), Inches(4.9), Inches(10.3), Inches(0.8), [
        {'text': 'Zero-cost, self-contained deployment: static assets and uploaded media are served '
                 'directly by the Python server. No CDN, no cloud storage, no external dependencies.',
         'sz': 11, 'color': C.T2, 'sp': 6},
        {'text': 'Request flow: Browser calls REST endpoints (/api/...) with a JWT bearer token '
                 '-> server validates, reads/writes SQLite -> JSON response returned.',
         'sz': 11, 'color': C.T2, 'sp': 4},
    ])

    white_footer(s3, 3)

    # ===================================================================
    #  SLIDE 4 — CORE FEATURES (White bg, reference slide 4 style)
    # ===================================================================
    s4 = prs.slides.add_slide(blank)
    set_bg(s4, C.WHITE)
    section_head(s4, "WHAT MEMBERS CAN DO TODAY", "Core Features")

    features = [
        ("Community Feed & Posts", "Text + photo posts, likes, threaded comments"),
        ("Direct Messaging", "Real-time 1:1 chat with read/delivered status"),
        ("H Rides Booking", "GPS canvas, rider-captain chat, OTP ride start"),
        ("Admin Control Center", "User management, moderation, broadcasts"),
        ("Notifications", "Multi-type alert center with unread badge counts"),
        ("Events & Network", "RSVP events, connection requests, global search"),
    ]

    for i, (title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + col * Inches(6.3)
        y = Inches(1.7) + row * Inches(1.55)

        # Navy card with circle icon
        rect(s4, x, y, Inches(5.9), Inches(1.35), C.NAVY)
        circ(s4, x + Inches(0.25), y + Inches(0.3), Inches(0.55), C.NAVY_MID)
        tx(s4, x + Inches(0.25), y + Inches(0.35), Inches(0.55), Inches(0.4),
           str(i+1), sz=16, color=C.W1, bold=True, align=PP_ALIGN.CENTER)
        tx(s4, x + Inches(1.0), y + Inches(0.2), Inches(4.6), Inches(0.3),
           title, sz=15, color=C.W1, bold=True)
        tx(s4, x + Inches(1.0), y + Inches(0.6), Inches(4.6), Inches(0.5),
           desc, sz=11, color=C.W2)

    white_footer(s4, 4)

    # ===================================================================
    #  SLIDE 5 — PROCONNECT DEEP DIVE
    # ===================================================================
    s5 = prs.slides.add_slide(blank)
    set_bg(s5, C.WHITE)
    section_head(s5, "PORTAL DEEP DIVE", "ProConnect — Career Network")

    pc_features = [
        ("JWT Authentication", "Secure login with Quick Demo mode for live investor demos"),
        ("Interactive Feed", "Create posts with photos, like, comment, and share in real-time"),
        ("Messaging Engine", "1:1 direct messages with typing indicators and read receipts"),
        ("Profile System", "Skills, education, projects, certifications, and cover photos"),
        ("Global Search", "Search across people, posts, hashtags, and events instantly"),
        ("Analytics Dashboard", "Profile views, post impressions, and engagement metrics"),
        ("Dark Mode", "CSS variable-based theme toggle with localStorage persistence"),
        ("Admin Panel", "Content moderation, reports queue, user management, system health"),
    ]

    for i, (title, desc) in enumerate(pc_features):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + col * Inches(6.3)
        y = Inches(1.65) + row * Inches(1.25)
        accent_dot(s5, x, y + Inches(0.05), Inches(0.5))
        tx(s5, x + Inches(0.2), y, Inches(5.7), Inches(0.26),
           title, sz=13, color=C.T1, bold=True)
        tx(s5, x + Inches(0.2), y + Inches(0.3), Inches(5.7), Inches(0.4),
           desc, sz=11, color=C.T2)

    white_footer(s5, 5)

    # ===================================================================
    #  SLIDE 6 — H RIDES DEEP DIVE
    # ===================================================================
    s6 = prs.slides.add_slide(blank)
    set_bg(s6, C.WHITE)
    section_head(s6, "PORTAL DEEP DIVE", "H Rides — Transit Platform")

    hr_features = [
        ("Live GPS Canvas", "HTML5 Canvas-rendered route simulation with animated vehicle tracking"),
        ("Dual Interface", "Switch between Rider (book rides) and Captain (accept & dispatch)"),
        ("In-App Chat", "Real-time rider-captain communication during active rides"),
        ("OTP Verification", "Secure ride-start confirmation via one-time passcode"),
        ("Ride State Machine", "Lifecycle: Pending > Accepted > In Progress > Completed"),
        ("Earnings Dashboard", "Captain partners track earnings, trip history, and ratings"),
    ]

    for i, (title, desc) in enumerate(hr_features):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + col * Inches(6.3)
        y = Inches(1.65) + row * Inches(1.5)
        accent_dot(s6, x, y + Inches(0.05), Inches(0.5))
        tx(s6, x + Inches(0.2), y, Inches(5.7), Inches(0.26),
           title, sz=13, color=C.T1, bold=True)
        tx(s6, x + Inches(0.2), y + Inches(0.3), Inches(5.7), Inches(0.4),
           desc, sz=11, color=C.T2)

    white_footer(s6, 6)

    # ===================================================================
    #  SLIDE 7 — QUALITY & TESTING (reference slide 8 style)
    # ===================================================================
    s7 = prs.slides.add_slide(blank)
    set_bg(s7, C.WHITE)
    section_head(s7, "AUTOMATED REGRESSION COVERAGE", "Testing & Quality Metrics")

    # Top metric boxes
    metrics = [
        ("51", "Automated Tests"),
        ("100%", "Pass Rate"),
        ("5", "Test Suites"),
        ("0", "Errors / Failures"),
    ]

    for i, (val, label) in enumerate(metrics):
        x = Inches(0.6) + i * Inches(3.15)
        rect(s7, x, Inches(1.7), Inches(2.9), Inches(1.5), C.PANEL)
        tx(s7, x, Inches(1.85), Inches(2.9), Inches(0.6),
           val, sz=36, color=C.NAVY, bold=True, align=PP_ALIGN.CENTER)
        tx(s7, x, Inches(2.5), Inches(2.9), Inches(0.3),
           label, sz=10, color=C.T3, bold=True, align=PP_ALIGN.CENTER)

    # Test suites list
    tx(s7, Inches(0.6), Inches(3.55), Inches(8), Inches(0.3),
       "SUITES EXECUTED", sz=12, color=C.NAVY, bold=True)

    suites = [
        "Hashtag & Product Upgrade -- 10/10 passed",
        "Profile Photo Lifecycle -- 15/15 passed",
        "Final Engineering Commandment -- 10/10 passed",
        "Master Presentation Suite -- 16/16 passed",
    ]

    for i, suite in enumerate(suites):
        y = Inches(4.1) + i * Inches(0.7)
        rnd(s7, Inches(0.6), y, Inches(12.1), Inches(0.55), C.PANEL)
        circ(s7, Inches(0.85), y + Inches(0.1), Inches(0.35), C.TEAL)
        tx(s7, Inches(0.85), y + Inches(0.12), Inches(0.35), Inches(0.3),
           "P", sz=11, color=C.WHITE, bold=True, align=PP_ALIGN.CENTER)  # "P" for pass
        tx(s7, Inches(1.4), y + Inches(0.1), Inches(11.0), Inches(0.35),
           suite, sz=12, color=C.T1)

    white_footer(s7, 7)

    # ===================================================================
    #  SLIDE 8 — COMPETITIVE EDGE (White bg)
    # ===================================================================
    s8 = prs.slides.add_slide(blank)
    set_bg(s8, C.WHITE)
    section_head(s8, "WHY WE WIN", "Competitive Advantage")

    rows_data = [
        ["FEATURE", "JOINING HANDS", "LINKEDIN", "RAPIDO"],
        ["Unified SSO", "Multi-Portal SSO", "Isolated", "Isolated"],
        ["Verification", "100% Hand-Verified", "Open Signups", "Phone Only"],
        ["Advertising", "Zero Ads", "Heavy Ad Load", "Sponsored"],
        ["Page Weight", "450KB SPA @ 60fps", "15MB+ Bundle", "8MB+ Native"],
        ["Privacy", "Zero Trackers", "Data Monetized", "Limited"],
        ["Ownership", "Community-First", "Corporate", "VC-Backed"],
    ]

    tbl_shape = s8.shapes.add_table(len(rows_data), 4, Inches(0.6), Inches(1.65), Inches(12.1), Inches(4.8))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.5)
    tbl.columns[1].width = Inches(3.8)
    tbl.columns[2].width = Inches(2.9)
    tbl.columns[3].width = Inches(2.9)

    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            p = cell.text_frame.paragraphs[0]
            p.text = val; p.font.name = C.FONT
            p.font.size = Pt(12 if ri == 0 else 11)
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT

            if ri == 0:
                cell.fill.fore_color.rgb = C.NAVY
                p.font.color.rgb = C.W1; p.font.bold = True
            elif ci == 1:
                cell.fill.fore_color.rgb = RGBColor(235, 245, 240)
                p.font.color.rgb = C.TEAL; p.font.bold = True
            elif ri % 2 == 0:
                cell.fill.fore_color.rgb = C.PANEL
                p.font.color.rgb = C.T2
            else:
                cell.fill.fore_color.rgb = C.WHITE
                p.font.color.rgb = C.T2

    white_footer(s8, 8)

    # ===================================================================
    #  SLIDE 9 — ROADMAP (White bg)
    # ===================================================================
    s9 = prs.slides.add_slide(blank)
    set_bg(s9, C.WHITE)
    section_head(s9, "WHAT'S NEXT", "Roadmap & Future Scope")

    milestones = [
        ("Q4 2026", "Platform Stable", "Production deployment, 51-test regression, admin health monitoring", C.TEAL),
        ("Q1 2027", "Welfare Launch", "Youth empowerment dashboard, donation ledger, event RSVP", C.SAFFRON),
        ("Q2 2027", "ML Career Match", "Skills-to-opportunity engine, content recommendations", C.BLUE),
        ("Q3 2027", "Encrypted Messaging", "WebRTC peer-to-peer, zero-knowledge storage, key rotation", C.NAVY),
    ]

    for i, (q, title, desc, clr) in enumerate(milestones):
        x = Inches(0.6) + i * Inches(3.15)
        rect(s9, x, Inches(1.7), Inches(2.95), Inches(3.5), C.NAVY)
        # Quarter badge
        rnd(s9, x + Inches(0.15), Inches(1.85), Inches(1.1), Inches(0.35), C.NAVY_MID)
        tx(s9, x + Inches(0.15), Inches(1.87), Inches(1.1), Inches(0.3),
           q, sz=10, color=C.LAVENDER, bold=True, align=PP_ALIGN.CENTER)
        # Title
        tx(s9, x + Inches(0.2), Inches(2.4), Inches(2.55), Inches(0.4),
           title, sz=16, color=C.W1, bold=True)
        # Description
        tx(s9, x + Inches(0.2), Inches(2.9), Inches(2.55), Inches(2.0),
           desc, sz=11, color=C.W2)

    white_footer(s9, 9)

    # ===================================================================
    #  SLIDE 10 — CLOSING (Navy bg, matching reference slide 9)
    # ===================================================================
    s10 = prs.slides.add_slide(blank)
    set_bg(s10, C.NAVY)

    tx(s10, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.35),
       "THANK YOU", sz=13, color=C.LAVENDER, bold=True, align=PP_ALIGN.CENTER)

    tx(s10, Inches(0.6), Inches(2.5), Inches(12.1), Inches(1.0),
       "Together Forever", sz=44, color=C.W1, bold=True, italic=True, align=PP_ALIGN.CENTER)

    rect(s10, Inches(5.5), Inches(3.7), Inches(2.3), Inches(0.015), C.LAVENDER)

    tx(s10, Inches(0.6), Inches(4.1), Inches(12.1), Inches(0.5),
       "The Group of Joining Hands", sz=18, color=C.W2, align=PP_ALIGN.CENTER)

    tx(s10, Inches(0.6), Inches(4.8), Inches(12.1), Inches(0.5),
       "Bengaluru, Karnataka, India  |  2026", sz=12, color=C.W3, align=PP_ALIGN.CENTER)

    # Checklist items (reference slide 9 style)
    checklist = [
        "Platform is production-ready with 51 automated tests at 100% pass rate",
        "Zero external dependencies: no cloud storage, no CDN, no third-party trackers",
        "Single-command deployment: python server.py starts the entire ecosystem",
        "Community-first model: members own their data, not advertisers",
    ]

    for i, item in enumerate(checklist):
        y = Inches(5.4) + i * Inches(0.38)
        circ(s10, Inches(3.5), y + Inches(0.05), Inches(0.2), C.LAVENDER)
        tx(s10, Inches(3.85), y, Inches(6.0), Inches(0.35),
           item, sz=10, color=C.W2)

    navy_footer(s10, 10)

    # --- SAVE ---
    out = "Joining_Hands_Investor_Deck.pptx"
    prs.save(out)
    print(f"[OK] Generated: {out} (10 slides)")


if __name__ == '__main__':
    build()
